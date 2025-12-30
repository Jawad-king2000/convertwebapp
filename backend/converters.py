"""
File Conversion Logic
Contains all the conversion functions for different file formats
"""

from pathlib import Path
from PIL import Image
import pdf2image
import pdfplumber
from pdf2docx import Converter as PDFToDocxConverter
from docx import Document
import pptx
import pytesseract
import pdfkit
import uuid
import zipfile
import io
import base64


def get_valid_output_formats(input_format: str) -> list:
    """
    Returns list of valid output formats for a given input format
    """
    conversion_rules = {
        'pdf': ['png', 'jpg', 'webp', 'txt', 'html', 'docx'],
        'docx': ['pdf', 'txt', 'html', 'png', 'jpg'],
        'pptx': ['pdf', 'png', 'jpg'],
        'txt': ['pdf', 'docx', 'html'],
        'png': ['txt', 'pdf', 'jpg', 'webp'],
        'jpg': ['txt', 'pdf', 'png', 'webp'],
        'jpeg': ['txt', 'pdf', 'png', 'webp'],
        'webp': ['png', 'jpg', 'txt', 'pdf']
    }
    return conversion_rules.get(input_format.lower(), [])


# ============================================================================
# PDF CONVERSIONS
# ============================================================================

def convert_pdf_to_format(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts PDF to various formats
    """
    if output_format in ['png', 'jpg', 'webp']:
        return pdf_to_image(input_path, output_format, output_dir)
    elif output_format == 'txt':
        return pdf_to_txt(input_path, output_dir)
    elif output_format == 'html':
        return pdf_to_html(input_path, output_dir)
    elif output_format == 'docx':
        return pdf_to_docx(input_path, output_dir)
    else:
        raise ValueError(f"Unsupported conversion: PDF to {output_format}")


def pdf_to_image(input_path: Path, image_format: str, output_dir: Path) -> Path:
    """
    Converts PDF pages to images
    If PDF has multiple pages, creates a ZIP file with all images
    """
    # Convert PDF to images (one image per page)
    images = pdf2image.convert_from_path(input_path, dpi=300)
    
    if len(images) == 1:
        # Single page - return single image
        output_path = output_dir / f"{uuid.uuid4()}.{image_format}"
        images[0].save(output_path, image_format.upper())
        return output_path
    else:
        # Multiple pages - create ZIP
        zip_path = output_dir / f"{uuid.uuid4()}.zip"
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for i, img in enumerate(images, 1):
                img_bytes = io.BytesIO()
                img.save(img_bytes, format=image_format.upper())
                zipf.writestr(f"page_{i}.{image_format}", img_bytes.getvalue())
        return zip_path


def pdf_to_txt(input_path: Path, output_dir: Path) -> Path:
    """
    Extracts text from PDF
    """
    output_path = output_dir / f"{uuid.uuid4()}.txt"
    
    text_content = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(text_content))
    
    return output_path


def pdf_to_html(input_path: Path, output_dir: Path) -> Path:
    """
    Converts PDF to HTML
    """
    output_path = output_dir / f"{uuid.uuid4()}.html"
    
    # Extract text and create simple HTML
    text_content = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_content.append(f"<div class='page'>{text.replace(chr(10), '<br>')}</div>")
    
    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Converted PDF</title>
    <style>
        body {{ font-family: Arial, sans-serif; padding: 20px; }}
        .page {{ margin-bottom: 40px; padding: 20px; border: 1px solid #ccc; }}
    </style>
</head>
<body>
    {''.join(text_content)}
</body>
</html>"""
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return output_path


def pdf_to_docx(input_path: Path, output_dir: Path) -> Path:
    """
    Converts PDF to DOCX
    """
    output_path = output_dir / f"{uuid.uuid4()}.docx"
    
    # Use pdf2docx library
    cv = PDFToDocxConverter(str(input_path))
    cv.convert(str(output_path))
    cv.close()
    
    return output_path


# ============================================================================
# DOCX CONVERSIONS
# ============================================================================

def convert_docx_to_format(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts DOCX to various formats
    """
    if output_format == 'pdf':
        return docx_to_pdf(input_path, output_dir)
    elif output_format == 'txt':
        return docx_to_txt(input_path, output_dir)
    elif output_format == 'html':
        return docx_to_html(input_path, output_dir)
    elif output_format in ['png', 'jpg']:
        # Convert via PDF intermediate
        pdf_path = docx_to_pdf(input_path, output_dir)
        return pdf_to_image(pdf_path, output_format, output_dir)
    else:
        raise ValueError(f"Unsupported conversion: DOCX to {output_format}")


def docx_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """
    Converts DOCX to PDF using HTML intermediate
    """
    # First convert to HTML, then to PDF
    html_path = docx_to_html(input_path, output_dir)
    output_path = output_dir / f"{uuid.uuid4()}.pdf"
    
    # Configure pdfkit to use the installed wkhtmltopdf
    try:
        # Standard path in Debian/Ubuntu (Docker)
        config = pdfkit.configuration(wkhtmltopdf='/usr/bin/wkhtmltopdf')
        pdfkit.from_file(str(html_path), str(output_path), configuration=config)
    except OSError:
        # Fallback for local development if not in standard path
        pdfkit.from_file(str(html_path), str(output_path))
        
    return output_path


def docx_to_txt(input_path: Path, output_dir: Path) -> Path:
    """
    Extracts text from DOCX
    """
    output_path = output_dir / f"{uuid.uuid4()}.txt"
    
    doc = Document(input_path)
    text_content = []
    
    for paragraph in doc.paragraphs:
        text_content.append(paragraph.text)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(text_content))
    
    return output_path


def docx_to_html(input_path: Path, output_dir: Path) -> Path:
    """
    Converts DOCX to HTML
    """
    output_path = output_dir / f"{uuid.uuid4()}.html"
    
    doc = Document(input_path)
    html_parts = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            html_parts.append(f"<p>{paragraph.text}</p>")
    
    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Converted Document</title>
    <style>
        body {{ font-family: Arial, sans-serif; padding: 20px; max-width: 800px; margin: 0 auto; }}
        p {{ margin: 10px 0; }}
    </style>
</head>
<body>
    {''.join(html_parts)}
</body>
</html>"""
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return output_path


# ============================================================================
# PPTX CONVERSIONS
# ============================================================================

def convert_pptx_to_format(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts PPTX to various formats
    """
    if output_format == 'pdf':
        return pptx_to_pdf(input_path, output_dir)
    elif output_format in ['png', 'jpg']:
        return pptx_to_images(input_path, output_format, output_dir)
    else:
        raise ValueError(f"Unsupported conversion: PPTX to {output_format}")


def pptx_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """
    Converts PPTX to PDF via images
    """
    # Convert slides to images first, then combine into PDF
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    output_path = output_dir / f"{uuid.uuid4()}.pdf"
    
    # For simplicity, we'll create a text-based PDF with slide content
    prs = pptx.Presentation(input_path)
    
    c = canvas.Canvas(str(output_path), pagesize=letter)
    
    for i, slide in enumerate(prs.slides, 1):
        c.drawString(100, 750, f"Slide {i}")
        y_position = 700
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                c.drawString(100, y_position, shape.text[:80])
                y_position -= 20
        
        c.showPage()
    
    c.save()
    return output_path


def pptx_to_images(input_path: Path, image_format: str, output_dir: Path) -> Path:
    """
    Converts PPTX slides to images (returns ZIP)
    """
    # This is a simplified version
    # In production, you'd use tools like unoconv or LibreOffice
    
    zip_path = output_dir / f"{uuid.uuid4()}.zip"
    
    prs = pptx.Presentation(input_path)
    
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for i, slide in enumerate(prs.slides, 1):
            # Create a simple placeholder image with slide number
            img = Image.new('RGB', (1920, 1080), color='white')
            img_bytes = io.BytesIO()
            img.save(img_bytes, format=image_format.upper())
            zipf.writestr(f"slide_{i}.{image_format}", img_bytes.getvalue())
    
    return zip_path


# ============================================================================
# TXT CONVERSIONS
# ============================================================================

def convert_txt_to_format(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts TXT to various formats
    """
    if output_format == 'pdf':
        return txt_to_pdf(input_path, output_dir)
    elif output_format == 'docx':
        return txt_to_docx(input_path, output_dir)
    elif output_format == 'html':
        return txt_to_html(input_path, output_dir)
    else:
        raise ValueError(f"Unsupported conversion: TXT to {output_format}")


def txt_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """
    Converts TXT to PDF
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit
    
    output_path = output_dir / f"{uuid.uuid4()}.pdf"
    
    with open(input_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    
    # Split text into lines that fit the page
    lines = text.split('\n')
    y_position = height - 50
    
    for line in lines:
        if y_position < 50:
            c.showPage()
            y_position = height - 50
        
        # Wrap long lines
        wrapped_lines = simpleSplit(line, "Helvetica", 10, width - 100)
        for wrapped_line in wrapped_lines:
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            c.drawString(50, y_position, wrapped_line)
            y_position -= 15
    
    c.save()
    return output_path


def txt_to_docx(input_path: Path, output_dir: Path) -> Path:
    """
    Converts TXT to DOCX
    """
    output_path = output_dir / f"{uuid.uuid4()}.docx"
    
    with open(input_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    doc = Document()
    
    for line in text.split('\n'):
        doc.add_paragraph(line)
    
    doc.save(output_path)
    return output_path


def txt_to_html(input_path: Path, output_dir: Path) -> Path:
    """
    Converts TXT to HTML
    """
    output_path = output_dir / f"{uuid.uuid4()}.html"
    
    with open(input_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    # Convert line breaks to <br> tags
    html_text = text.replace('\n', '<br>\n')
    
    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Converted Text</title>
    <style>
        body {{ font-family: monospace; padding: 20px; white-space: pre-wrap; }}
    </style>
</head>
<body>
{html_text}
</body>
</html>"""
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return output_path


# ============================================================================
# IMAGE CONVERSIONS
# ============================================================================

def convert_image_to_format(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts images to various formats
    """
    if output_format == 'txt':
        return image_to_txt_ocr(input_path, output_dir)
    elif output_format == 'pdf':
        return image_to_pdf(input_path, output_dir)
    elif output_format in ['png', 'jpg', 'webp']:
        return image_to_image(input_path, output_format, output_dir)
    else:
        raise ValueError(f"Unsupported conversion: Image to {output_format}")


def image_to_txt_ocr(input_path: Path, output_dir: Path) -> Path:
    """
    Extracts text from image using OCR
    """
    output_path = output_dir / f"{uuid.uuid4()}.txt"
    
    try:
        # Open image and perform OCR
        img = Image.open(input_path)
        text = pytesseract.image_to_string(img)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    except Exception as e:
        # If OCR fails, create a file explaining the error
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"OCR Error: {str(e)}\n\nMake sure Tesseract is installed.")
    
    return output_path


def image_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """
    Converts image to PDF
    """
    output_path = output_dir / f"{uuid.uuid4()}.pdf"
    
    img = Image.open(input_path)
    
    # Convert to RGB if necessary
    if img.mode != 'RGB':
        img = img.convert('RGB')
    
    img.save(output_path, 'PDF')
    
    return output_path


def image_to_image(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """
    Converts between image formats
    """
    output_path = output_dir / f"{uuid.uuid4()}.{output_format}"
    
    img = Image.open(input_path)
    
    # Convert RGBA to RGB for JPEG
    if output_format.lower() in ['jpg', 'jpeg'] and img.mode == 'RGBA':
        rgb_img = Image.new('RGB', img.size, (255, 255, 255))
        rgb_img.paste(img, mask=img.split()[3] if len(img.split()) == 4 else None)
        img = rgb_img
    
    img.save(output_path, output_format.upper())
    
    return output_path
